#!/usr/bin/env python
# -*- coding: utf-8 -*-

import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from collections import Counter
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import numpy as np

# Create output directories
if not os.path.exists('output_images'):
    os.makedirs('output_images')

# Load AWS ID data
df = pd.read_csv('aws_id.csv')

# Basic data exploration
total_records = len(df)
unique_aws_ids = df['aws_id'].nunique()
unique_account_ids = df['account_id'].nunique()
duplicate_aws_ids = total_records - unique_aws_ids

# Find AWS IDs with multiple accounts
aws_id_counts = df['aws_id'].value_counts()
aws_ids_with_multiple_accounts = aws_id_counts[aws_id_counts > 1].index.tolist()
multi_account_data = df[df['aws_id'].isin(aws_ids_with_multiple_accounts)].sort_values('aws_id')

# Count AWS IDs by prefix and numerical part
aws_id_prefixes = df['aws_id'].apply(lambda x: x[:3])  # 'AWS' prefix
aws_id_numbers = df['aws_id'].apply(lambda x: x[3:-1])  # Numerical part
aws_id_numbers_int = pd.to_numeric(aws_id_numbers)
aws_id_suffix = df['aws_id'].apply(lambda x: x[-1])  # 'S' suffix

prefix_counts = aws_id_prefixes.value_counts().reset_index()
prefix_counts.columns = ['Prefix', 'Count']

suffix_counts = aws_id_suffix.value_counts().reset_index()
suffix_counts.columns = ['Suffix', 'Count']

# Create visualizations

# 1. Bar chart of AWS IDs with multiple accounts
plt.figure(figsize=(12, 6))
top_duplicate_aws_ids = aws_id_counts[aws_id_counts > 1].head(10)
plt.bar(top_duplicate_aws_ids.index, top_duplicate_aws_ids.values, color='skyblue')
plt.title('Top 10 AWS IDs with Multiple Accounts')
plt.xlabel('AWS ID')
plt.ylabel('Number of Accounts')
plt.xticks(rotation=90)
plt.tight_layout()
plt.savefig('output_images/aws_ids_with_multiple_accounts.png', dpi=300)
plt.close()

# 2. Histogram of account ID length
plt.figure(figsize=(10, 6))
account_id_length = df['account_id'].astype(str).apply(len)
sns.histplot(account_id_length, kde=True, bins=range(min(account_id_length), max(account_id_length) + 1))
plt.title('Distribution of Account ID Length')
plt.xlabel('Length of Account ID')
plt.ylabel('Frequency')
plt.tight_layout()
plt.savefig('output_images/account_id_length_distribution.png', dpi=300)
plt.close()

# 3. Distribution of AWS ID numerical values
plt.figure(figsize=(12, 6))
sns.histplot(aws_id_numbers_int, kde=True, bins=30)
plt.title('Distribution of AWS ID Numerical Values')
plt.xlabel('AWS ID Number')
plt.ylabel('Frequency')
plt.tight_layout()
plt.savefig('output_images/aws_id_number_distribution.png', dpi=300)
plt.close()

# 4. Pie chart for duplicate vs unique AWS IDs
plt.figure(figsize=(10, 6))
labels = ['Unique AWS IDs', 'Duplicate AWS IDs']
sizes = [unique_aws_ids, duplicate_aws_ids]
colors = ['lightblue', 'lightcoral']
plt.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%', startangle=90)
plt.axis('equal')
plt.title('Proportion of Unique vs Duplicate AWS IDs')
plt.tight_layout()
plt.savefig('output_images/unique_vs_duplicate_aws_ids.png', dpi=300)
plt.close()

# 5. Number of accounts per AWS ID visualization
plt.figure(figsize=(10, 6))
accounts_per_id = df['aws_id'].value_counts()
sns.histplot(accounts_per_id.values, kde=True, bins=range(1, max(accounts_per_id.values) + 2))
plt.title('Number of Accounts per AWS ID')
plt.xlabel('Number of Accounts')
plt.ylabel('Count of AWS IDs')
plt.tight_layout()
plt.savefig('output_images/accounts_per_aws_id.png', dpi=300)
plt.close()

# 6. Account ID first digit distribution
plt.figure(figsize=(10, 6))
first_digits = df['account_id'].astype(str).str[0].value_counts().sort_index()
sns.barplot(x=first_digits.index, y=first_digits.values)
plt.title('First Digit Distribution of Account IDs')
plt.xlabel('First Digit')
plt.ylabel('Frequency')
plt.tight_layout()
plt.savefig('output_images/account_id_first_digit.png', dpi=300)
plt.close()

# Create PowerPoint presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AWS ID Analysis"
subtitle.text = f"Total Records: {total_records}, Unique AWS IDs: {unique_aws_ids}, Unique Account IDs: {unique_account_ids}"

# Function to add a slide with an image
def add_image_slide(title_text, img_path, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = title_text
    
    # Add image
    left = Inches(1.5)
    top = Inches(2)
    slide.shapes.add_picture(img_path, left, top, width=Inches(7))
    
    # Add subtitle if provided
    if subtitle_text:
        text_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.5))
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = subtitle_text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)

# Summary slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Data Summary"

# Add summary text
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = f"• Total Records: {total_records}"
p = tf.add_paragraph()
p.text = f"• Unique AWS IDs: {unique_aws_ids}"
p = tf.add_paragraph()
p.text = f"• Unique Account IDs: {unique_account_ids}"
p = tf.add_paragraph()
p.text = f"• AWS IDs with Multiple Accounts: {len(aws_ids_with_multiple_accounts)}"
p = tf.add_paragraph()
p.text = f"• Most common AWS ID suffix: {suffix_counts.iloc[0]['Suffix']} ({suffix_counts.iloc[0]['Count']} occurrences)"
p = tf.add_paragraph()
p.text = f"• Account ID length: Typically {account_id_length.mode()[0]} digits"

# Add image slides
add_image_slide("Top AWS IDs with Multiple Accounts", "output_images/aws_ids_with_multiple_accounts.png",
                f"There are {len(aws_ids_with_multiple_accounts)} AWS IDs with multiple accounts")

add_image_slide("Distribution of Account ID Length", "output_images/account_id_length_distribution.png",
                "Most account IDs have a consistent length")

add_image_slide("Distribution of AWS ID Numerical Values", "output_images/aws_id_number_distribution.png",
                "Analysis of the numerical portion of AWS IDs")

add_image_slide("Unique vs Duplicate AWS IDs", "output_images/unique_vs_duplicate_aws_ids.png",
                f"Unique AWS IDs: {unique_aws_ids}, Duplicate AWS IDs: {duplicate_aws_ids}")

add_image_slide("Number of Accounts per AWS ID", "output_images/accounts_per_aws_id.png",
                "Distribution showing how many accounts are associated with each AWS ID")

add_image_slide("First Digit Distribution of Account IDs", "output_images/account_id_first_digit.png",
                "Distribution of the first digit in account IDs")

# AWS IDs with Multiple Accounts slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "AWS IDs with Multiple Accounts"

# Add table with data
if multi_account_data.shape[0] > 0:
    rows = min(11, multi_account_data.shape[0] + 1)  # +1 for header
    cols = 2
    
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(0.8 * rows)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set header
    table.cell(0, 0).text = "AWS ID"
    table.cell(0, 1).text = "Account ID"
    
    # Fill data
    for i in range(min(10, multi_account_data.shape[0])):
        table.cell(i + 1, 0).text = str(multi_account_data.iloc[i]['aws_id'])
        table.cell(i + 1, 1).text = str(multi_account_data.iloc[i]['account_id'])
else:
    body_shape = slide.placeholders[1]
    body_shape.text_frame.text = "No AWS IDs with multiple accounts found"

# Conclusion slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Key Findings:"
p = tf.add_paragraph()
p.text = f"• Dataset contains {total_records} total AWS ID to account ID mappings"
p = tf.add_paragraph()
p.text = f"• {len(aws_ids_with_multiple_accounts)} AWS IDs are associated with multiple accounts"
p = tf.add_paragraph()
p.text = f"• Account ID length is consistently {account_id_length.mode()[0]} digits for most records"
p = tf.add_paragraph()
p.text = "• The AWS ID naming pattern follows a consistent format (AWS followed by 5 digits and S)"
p = tf.add_paragraph()
p.text = f"• There are {unique_account_ids} unique account IDs in the dataset"

# Save presentation
prs.save('AWS_ID_Analysis.pptx')

# Create PDF report
class PDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 12)
        self.cell(0, 10, 'AWS ID Analysis Report', 0, 1, 'C')
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
        
    def chapter_title(self, title):
        self.set_font('Arial', 'B', 12)
        self.cell(0, 10, title, 0, 1, 'L')
        self.ln(4)
        
    def chapter_body(self, body):
        self.set_font('Arial', '', 10)
        self.multi_cell(0, 5, body)
        self.ln()
        
    def add_image(self, img_path, w=190):
        self.image(img_path, x=10, w=w)
        self.ln(5)

pdf = PDF()
pdf.set_auto_page_break(auto=True, margin=15)
pdf.add_page()

# Introduction
pdf.chapter_title('1. Introduction')
pdf.chapter_body(f'This report analyzes AWS ID and account ID data from the aws_id.csv file. '
                f'The dataset contains {total_records} records with {unique_aws_ids} unique AWS IDs '
                f'and {unique_account_ids} unique account IDs.')

# Data Overview
pdf.add_page()
pdf.chapter_title('2. Data Overview')
pdf.chapter_body(f'Total Records: {total_records}\n'
                f'Unique AWS IDs: {unique_aws_ids}\n'
                f'Unique Account IDs: {unique_account_ids}\n'
                f'AWS IDs with Multiple Accounts: {len(aws_ids_with_multiple_accounts)}\n'
                f'Account ID Length: Most are {account_id_length.mode()[0]} digits\n\n'
                f'AWS ID Format: All IDs follow the pattern "AWS" + 5 digits + "S"')

# Visualizations
pdf.add_page()
pdf.chapter_title('3. Visualizations')
pdf.chapter_body('3.1 AWS IDs with Multiple Accounts')
pdf.add_image('output_images/aws_ids_with_multiple_accounts.png')
pdf.chapter_body(f'There are {len(aws_ids_with_multiple_accounts)} AWS IDs with multiple accounts.')

pdf.add_page()
pdf.chapter_body('3.2 Distribution of Account ID Length')
pdf.add_image('output_images/account_id_length_distribution.png')
pdf.chapter_body(f'Most account IDs have a length of {account_id_length.mode()[0]} digits.')

pdf.add_page()
pdf.chapter_body('3.3 Distribution of AWS ID Numerical Values')
pdf.add_image('output_images/aws_id_number_distribution.png')
pdf.chapter_body('The numerical portion of AWS IDs shows specific patterns in their distribution.')

pdf.add_page()
pdf.chapter_body('3.4 Unique vs Duplicate AWS IDs')
pdf.add_image('output_images/unique_vs_duplicate_aws_ids.png')
pdf.chapter_body(f'Unique AWS IDs: {unique_aws_ids}, Duplicate AWS IDs: {duplicate_aws_ids}')

pdf.add_page()
pdf.chapter_body('3.5 Number of Accounts per AWS ID')
pdf.add_image('output_images/accounts_per_aws_id.png')
pdf.chapter_body('Distribution showing how many accounts are associated with each AWS ID.')

pdf.add_page()
pdf.chapter_body('3.6 First Digit Distribution of Account IDs')
pdf.add_image('output_images/account_id_first_digit.png')
pdf.chapter_body('Distribution of the first digit in account IDs.')

# AWS IDs with Multiple Accounts
pdf.add_page()
pdf.chapter_title('4. AWS IDs with Multiple Accounts')
if multi_account_data.shape[0] > 0:
    # Create a table-like structure for the first 10 entries
    top_multi_accounts = multi_account_data.head(10)
    pdf.set_font('Arial', 'B', 10)
    pdf.cell(90, 7, 'AWS ID', 1, 0, 'C')
    pdf.cell(90, 7, 'Account ID', 1, 1, 'C')
    
    pdf.set_font('Arial', '', 10)
    for _, row in top_multi_accounts.iterrows():
        pdf.cell(90, 7, str(row['aws_id']), 1, 0, 'L')
        pdf.cell(90, 7, str(row['account_id']), 1, 1, 'L')
else:
    pdf.chapter_body('No AWS IDs with multiple accounts found.')

# Conclusion
pdf.add_page()
pdf.chapter_title('5. Conclusion')
pdf.chapter_body(f'Key Findings:\n\n'
                f'• Dataset contains {total_records} total AWS ID to account ID mappings\n'
                f'• {len(aws_ids_with_multiple_accounts)} AWS IDs are associated with multiple accounts\n'
                f'• Account ID length is consistently {account_id_length.mode()[0]} digits for most records\n'
                f'• The AWS ID naming pattern follows a consistent format (AWS followed by 5 digits and S)\n'
                f'• There are {unique_account_ids} unique account IDs in the dataset')

# Save PDF
pdf.output('AWS_ID_Analysis.pdf')

print(f"Analysis complete. PowerPoint saved as AWS_ID_Analysis.pptx")
print(f"PDF report saved as AWS_ID_Analysis.pdf")
print(f"Generated {len(os.listdir('output_images'))} visualization images in the output_images folder")